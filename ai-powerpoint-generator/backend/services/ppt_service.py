from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional
import io
import os
from datetime import datetime

class PowerPointService:
    """Service for generating PowerPoint presentations"""

    def __init__(self):
        self.presentation = None

    def create_presentation(self, slides_data: List[Dict], images: Dict[int, io.BytesIO], title: str) -> str:
        """
        Create a PowerPoint presentation with content and images

        Args:
            slides_data: List of slide dictionaries with title, content, and keywords
            images: Dictionary mapping slide numbers to image BytesIO objects
            title: Presentation title

        Returns:
            Path to the generated PowerPoint file
        """
        self.presentation = Presentation()

        # Set slide dimensions (16:9)
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(5.625)

        for slide_data in slides_data:
            slide_num = slide_data.get("slide_number", 1)
            slide_title = slide_data.get("title", "")
            slide_content = slide_data.get("content", [])
            image = images.get(slide_num)

            if slide_num == 1:
                # Title slide
                self._create_title_slide(slide_title, slide_content)
            else:
                # Content slide
                self._create_content_slide(slide_title, slide_content, image)

        # Save the presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        output_path = os.path.join("../output", filename)

        # Ensure output directory exists
        os.makedirs("../output", exist_ok=True)

        self.presentation.save(output_path)
        return filename

    def _create_title_slide(self, title: str, subtitle_points: List[str]):
        """Create an attractive title slide"""
        # Use blank layout and add shapes manually for better control
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 121)  # Professional blue

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(9),
            Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        # Format title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add subtitle if available
        if subtitle_points:
            subtitle_text = " • ".join(subtitle_points[:3])
            subtitle_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(3.5),
                Inches(8),
                Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.word_wrap = True

            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = RGBColor(220, 220, 220)

    def _create_content_slide(self, title: str, content: List[str], image: Optional[io.BytesIO]):
        """Create a content slide with title, bullet points, and optional image"""
        # Use blank layout for full control
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Determine layout based on image presence
        if image:
            # Two-column layout: content on left, image on right
            content_left = Inches(0.5)
            content_width = Inches(4.5)
            image_left = Inches(5.5)
            image_width = Inches(4)
        else:
            # Full-width content
            content_left = Inches(0.5)
            content_width = Inches(9)

        # Add title with colored background
        title_box = slide.shapes.add_textbox(
            Inches(0),
            Inches(0),
            Inches(10),
            Inches(0.8)
        )

        # Title background
        title_fill = title_box.fill
        title_fill.solid()
        title_fill.fore_color.rgb = RGBColor(31, 78, 121)

        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.margin_left = Inches(0.3)
        title_frame.margin_top = Inches(0.15)

        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add content
        content_box = slide.shapes.add_textbox(
            content_left,
            Inches(1.2),
            content_width,
            Inches(4)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = MSO_ANCHOR.TOP

        for i, point in enumerate(content):
            if i > 0:
                content_frame.add_paragraph()

            para = content_frame.paragraphs[i]
            para.text = point
            para.level = 0
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(50, 50, 50)
            para.space_before = Pt(12)

            # Add bullet
            para.bullet = True

        # Add image if available
        if image:
            try:
                # Reset BytesIO position
                image.seek(0)

                # Add image to slide
                pic = slide.shapes.add_picture(
                    image,
                    image_left,
                    Inches(1.2),
                    width=image_width
                )

                # Adjust height to maintain aspect ratio
                if pic.height > Inches(4):
                    pic.height = Inches(4)

            except Exception as e:
                print(f"Error adding image to slide: {e}")

    def _add_conclusion_slide(self, title: str = "Thank You"):
        """Add a conclusion slide"""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 121)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
